"""
Office Document Extractor Module

This module provides functionality to extract text content from various Microsoft Office
document formats including DOCX, Excel (XLS, XLSX, XLSM), and PowerPoint (PPT, PPTX).
"""

import io
import logging
import pandas as pd
import requests
from typing import Union, List
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import xlrd
from urllib.parse import urlparse

logger = logging.getLogger(__name__)

class OfficeExtractor:
    def __init__(self):
        """Initialize the OfficeExtractor."""
        pass

    def _download_file(self, url: str) -> bytes:
        """
        Download a file from a URL.

        Args:
            url (str): The URL of the file to download.

        Returns:
            bytes: The downloaded file content.
        """
        try:
            response = requests.get(url, verify=True)
            response.raise_for_status()
            return response.content
        except Exception as e:
            logger.error(f"Error downloading file from {url}: {str(e)}")
            raise

    def _is_url(self, source: str) -> bool:
        """Check if the source is a URL."""
        try:
            result = urlparse(source)
            return all([result.scheme, result.netloc])
        except ValueError:
            return False

    def extract_docx(self, source: Union[str, bytes]) -> str:
        """
        Extract text from a DOCX file.

        Args:
            source (Union[str, bytes]): File path, URL, or bytes content.

        Returns:
            str: Extracted text content.
        """
        try:
            if isinstance(source, str):
                if self._is_url(source):
                    content = self._download_file(source)
                    doc = Document(io.BytesIO(content))
                else:
                    doc = Document(source)
            else:
                doc = Document(io.BytesIO(source))

            # Extract text from paragraphs
            text = []
            for paragraph in doc.paragraphs:
                text.append(paragraph.text)

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text)
                    text.append(" | ".join(row_text))

            return "\n".join(text)
        except Exception as e:
            logger.error(f"Error extracting DOCX content: {str(e)}")
            raise

    def extract_excel(self, source: Union[str, bytes]) -> str:
        """
        Extract text from Excel files (XLS, XLSX, XLSM).

        Args:
            source (Union[str, bytes]): File path, URL, or bytes content.

        Returns:
            str: Extracted text content.
        """
        try:
            if isinstance(source, str):
                if self._is_url(source):
                    content = self._download_file(source)
                    if source.lower().endswith('.xls'):
                        # Handle old XLS format
                        workbook = xlrd.open_workbook(file_contents=content)
                        dfs = []
                        for sheet in workbook.sheets():
                            data = []
                            for row in range(sheet.nrows):
                                data.append([str(sheet.cell_value(row, col)) for col in range(sheet.ncols)])
                            dfs.append(pd.DataFrame(data[1:], columns=data[0]))
                    else:
                        # Handle XLSX/XLSM format
                        dfs = pd.read_excel(io.BytesIO(content), sheet_name=None)
                else:
                    if source.lower().endswith('.xls'):
                        workbook = xlrd.open_workbook(source)
                        dfs = []
                        for sheet in workbook.sheets():
                            data = []
                            for row in range(sheet.nrows):
                                data.append([str(sheet.cell_value(row, col)) for col in range(sheet.ncols)])
                            dfs.append(pd.DataFrame(data[1:], columns=data[0]))
                    else:
                        dfs = pd.read_excel(source, sheet_name=None)
            else:
                dfs = pd.read_excel(io.BytesIO(source), sheet_name=None)

            # Convert all sheets to text
            text = []
            for sheet_name, df in (dfs.items() if isinstance(dfs, dict) else enumerate(dfs)):
                text.append(f"\nSheet: {sheet_name}")
                text.append(df.to_string(index=False))

            return "\n".join(text)
        except Exception as e:
            logger.error(f"Error extracting Excel content: {str(e)}")
            raise

    def extract_pptx(self, source: Union[str, bytes]) -> str:
        """
        Extract text from PowerPoint files (PPTX).

        Args:
            source (Union[str, bytes]): File path, URL, or bytes content.

        Returns:
            str: Extracted text content.
        """
        try:
            if isinstance(source, str):
                if self._is_url(source):
                    content = self._download_file(source)
                    prs = Presentation(io.BytesIO(content))
                else:
                    prs = Presentation(source)
            else:
                prs = Presentation(io.BytesIO(source))

            text = []
            for slide_number, slide in enumerate(prs.slides, 1):
                slide_text = [f"\nSlide {slide_number}:"]
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            slide_text.append(shape.text)
                text.append("\n".join(slide_text))

            return "\n".join(text)
        except Exception as e:
            logger.error(f"Error extracting PowerPoint content: {str(e)}")
            raise

def main():
    """Test the OfficeExtractor class."""
    extractor = OfficeExtractor()
    
    # Test URLs
    test_urls = [
        "https://calibre-ebook.com/downloads/demos/demo.docx",
        "https://www.cmu.edu/blackboard/files/evaluate/tests-example.xls",
        "https://scholar.harvard.edu/files/torman_personal/files/samplepptx.pptx"
    ]
    
    for url in test_urls:
        try:
            if url.endswith('.docx'):
                content = extractor.extract_docx(url)
            elif url.endswith('.xls') or url.endswith('.xlsx'):
                content = extractor.extract_excel(url)
            elif url.endswith('.pptx'):
                content = extractor.extract_pptx(url)
                
            print(f"\nExtracted content from {url}:")
            print(content[:500] + "..." if len(content) > 500 else content)
        except Exception as e:
            print(f"Error processing {url}: {str(e)}")

if __name__ == "__main__":
    main() 