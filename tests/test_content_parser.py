"""
Test Content Parser Module

This module contains comprehensive tests for content extraction functionality,
including edge cases and error handling for various file types.
"""

import os
import unittest
import pytest
import tempfile
import requests
from unittest.mock import Mock, patch
from io import BytesIO
import pandas as pd

from podcastfy.utils.config import load_config
from podcastfy.content_parser.content_extractor import ContentExtractor
from podcastfy.content_parser.youtube_transcriber import YouTubeTranscriber
from podcastfy.content_parser.website_extractor import WebsiteExtractor
from podcastfy.content_parser.pdf_extractor import PDFExtractor
from podcastfy.content_parser.office_extractor import OfficeExtractor
from podcastfy.content_parser.text_extractor import TextExtractor

# Test data URLs
TEST_URLS = {
    'pdf': 'https://www.banyantree.in/jagdishpur/wp-content/uploads/2020/06/Panchatantra-.pdf',
    'docx': 'https://calibre-ebook.com/downloads/demos/demo.docx',
    'excel': 'https://file-examples.com/wp-content/storage/2017/02/file_example_XLS_10.xls',
    'pptx': 'https://scholar.harvard.edu/files/torman_personal/files/samplepptx.pptx',
    'website': 'https://example.com',
    'youtube': 'https://www.youtube.com/watch?v=dQw4w9WgXcQ',
    'text': 'https://example-files.online-convert.com/document/txt/example.txt',
    'text_iso': 'https://www.w3.org/TR/2003/REC-PNG-20031110/iso_8859-1.txt'
}

class TestContentParser(unittest.TestCase):
    @pytest.mark.skip(
        reason="IP getting blocked by YouTube when running from GitHub Actions"
    )
    def test_youtube_transcriber(self):
        """
        Test the YouTubeTranscriber class to ensure it correctly extracts and cleans transcripts from a YouTube video.
        """
        # Initialize YouTubeTranscriber
        transcriber = YouTubeTranscriber()

        # Test URL
        test_url = "https://www.youtube.com/watch?v=m3kJo5kEzTQ"

        # Extract transcript
        extracted_transcript = transcriber.extract_transcript(test_url)

        # Load expected transcript from youtube.txt file
        with open("./tests/data/mock/youtube.txt", "r") as f:
            expected_transcript = f.read()

        # Assert that the first 100 characters of the extracted transcript match the expected transcript
        self.assertEqual(
            extracted_transcript[:100].strip(), expected_transcript[:100].strip()
        )

    def test_website_extractor(self):
        """
        Test the WebsiteExtractor class to ensure it correctly extracts content from a website.
        """
        # Initialize WebsiteExtractor
        config = load_config()
        extractor = WebsiteExtractor()

        # Test URL
        test_url = "http://www.souzatharsis.com"

        # Extract content
        extracted_content = extractor.extract_content(test_url)
        print(extracted_content.strip())
        # Load expected content from website.md file
        with open("./tests/data/mock/website.md", "r") as f:
            expected_content = f.read()
        print(expected_content.strip())
        # Assert that the extracted content matches the expected content
        self.assertEqual(extracted_content.strip(), expected_content.strip())

    def test_pdf_extractor(self):
        """
        Test the PDFExtractor class to ensure it correctly extracts content from a PDF file.
        """
        # Initialize PDFExtractor
        extractor = PDFExtractor()

        # Path to the test PDF file
        pdf_path = "./tests/data/pdf/file.pdf"

        # Extract content from PDF
        extracted_content = extractor.extract_content(pdf_path)

        # Load expected content from file.txt
        with open("./tests/data/mock/file.txt", "r") as f:
            expected_content = f.read()

        # Assert that the first 500 characters of the extracted content match the expected content
        self.assertEqual(
            extracted_content[:500].strip(), expected_content[:500].strip()
        )

    @pytest.mark.skip(reason="Too expensive to be auto tested on Github Actions")
    def test_generate_topic_content(self):
        """Test generating content for a specific topic."""
        extractor = ContentExtractor()
        topic = "Latest news about OpenAI"

        # Generate content for the topic
        content = extractor.generate_topic_content(topic)

        # Verify the content
        self.assertIsNotNone(content)
        self.assertIsInstance(content, str)
        self.assertGreater(len(content), 100)  # Content should be substantial

        # Check if content is relevant to the topic
        lower_content = content.lower()
        self.assertTrue(
            any(term in lower_content for term in ["openai"]),
            "Generated content should be relevant to the topic",
        )

# New pytest-style tests for enhanced functionality
@pytest.fixture
def content_extractor():
    """Create a ContentExtractor instance for testing."""
    return ContentExtractor()

@pytest.fixture
def pdf_extractor():
    """Create a PDFExtractor instance for testing."""
    return PDFExtractor()

@pytest.fixture
def office_extractor():
    """Create an OfficeExtractor instance for testing."""
    return OfficeExtractor()

@pytest.fixture
def text_extractor():
    """Create a TextExtractor instance for testing."""
    return TextExtractor()

def create_mock_response(content: bytes, status_code: int = 200):
    """Create a mock requests.Response object."""
    mock_response = Mock(spec=requests.Response)
    mock_response.content = content
    mock_response.status_code = status_code
    mock_response.raise_for_status = Mock()
    if status_code >= 400:
        mock_response.raise_for_status.side_effect = requests.exceptions.HTTPError()
    return mock_response

class TestContentExtractor:
    def test_url_detection(self, content_extractor):
        """Test URL detection for various formats."""
        assert content_extractor.is_url('https://example.com')
        assert content_extractor.is_url('http://example.com')
        assert content_extractor.is_url('example.com')
        assert not content_extractor.is_url('path/to/file.pdf')
        assert not content_extractor.is_url('file.pdf')

    def test_file_extension_detection(self, content_extractor):
        """Test file extension detection."""
        assert content_extractor._get_file_extension('file.pdf') == 'pdf'
        assert content_extractor._get_file_extension('file.PDF') == 'pdf'
        assert content_extractor._get_file_extension('path/to/file.docx') == 'docx'
        assert content_extractor._get_file_extension('https://example.com/file.xlsx') == 'xlsx'
        assert content_extractor._get_file_extension('noextension') == ''

    @patch('requests.get')
    def test_web_pdf_extraction(self, mock_get, content_extractor):
        """Test PDF extraction from web URLs."""
        # Create a simple PDF content mock
        mock_pdf_content = b'%PDF-1.4\nsome content'
        mock_get.return_value = create_mock_response(mock_pdf_content)

        # Test successful extraction
        content = content_extractor.extract_content(TEST_URLS['pdf'])
        assert isinstance(content, str)
        mock_get.assert_called_once()

        # Test failed download
        mock_get.reset_mock()
        mock_get.return_value = create_mock_response(b'', status_code=404)
        with pytest.raises(requests.exceptions.HTTPError):
            content_extractor.extract_content(TEST_URLS['pdf'])

class TestOfficeExtractor:
    @patch('requests.get')
    def test_excel_extraction(self, mock_get, office_extractor):
        """Test Excel file extraction with various formats."""
        # Create mock Excel data
        df = pd.DataFrame({'A': [1, 2, 3], 'B': ['a', 'b', 'c']})
        excel_buffer = BytesIO()
        df.to_excel(excel_buffer, index=False)
        excel_content = excel_buffer.getvalue()

        # Test XLSX
        mock_get.return_value = create_mock_response(excel_content)
        content = office_extractor.extract_excel(TEST_URLS['excel'])
        assert isinstance(content, str)
        assert 'Sheet' in content
        
        # Test with merged cells
        df_merged = pd.DataFrame({'A': [1, 2, 3], 'B': ['a', 'b', 'c']})
        excel_buffer = BytesIO()
        with pd.ExcelWriter(excel_buffer, engine='openpyxl') as writer:
            df_merged.to_excel(writer, index=False)
            worksheet = writer.sheets['Sheet1']
            worksheet.merge_cells('A1:B1')
        mock_get.return_value = create_mock_response(excel_buffer.getvalue())
        content = office_extractor.extract_excel(TEST_URLS['excel'])
        assert isinstance(content, str)

    @patch('requests.get')
    def test_docx_extraction(self, mock_get, office_extractor):
        """Test DOCX extraction with various content types."""
        from docx import Document
        
        # Create a test DOCX with tables and text
        doc = Document()
        doc.add_paragraph('Test paragraph')
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = 'Cell 1'
        table.cell(0, 1).text = 'Cell 2'
        
        docx_buffer = BytesIO()
        doc.save(docx_buffer)
        mock_get.return_value = create_mock_response(docx_buffer.getvalue())
        
        content = office_extractor.extract_docx(TEST_URLS['docx'])
        assert isinstance(content, str)
        assert 'Test paragraph' in content
        assert 'Cell 1' in content
        assert 'Cell 2' in content

    @patch('requests.get')
    def test_pptx_extraction(self, mock_get, office_extractor):
        """Test PowerPoint extraction."""
        from pptx import Presentation
        
        # Create a test PPTX
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = 'Test Slide'
        
        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        mock_get.return_value = create_mock_response(pptx_buffer.getvalue())
        
        content = office_extractor.extract_pptx(TEST_URLS['pptx'])
        assert isinstance(content, str)
        assert 'Test Slide' in content

    def test_large_excel_file(self, office_extractor):
        """Test handling of large Excel files."""
        # Create a large Excel file
        large_df = pd.DataFrame({
            'A': range(10000),
            'B': ['test'] * 10000
        })
        excel_buffer = BytesIO()
        large_df.to_excel(excel_buffer, index=False)
        content = office_extractor.extract_excel(excel_buffer.getvalue())
        assert isinstance(content, str)
        assert len(content) > 1000

class TestTextExtractor:
    @patch('requests.get')
    def test_text_extraction(self, mock_get, text_extractor):
        """Test text file extraction with various encodings."""
        # Test UTF-8
        utf8_content = "Hello, World!".encode('utf-8')
        mock_get.return_value = create_mock_response(utf8_content)
        content = text_extractor.extract_content(TEST_URLS['text'])
        assert isinstance(content, str)
        assert content == "Hello, World!"

        # Test ISO-8859-1
        iso_content = "Hello, World!".encode('iso-8859-1')
        mock_get.return_value = create_mock_response(iso_content)
        content = text_extractor.extract_content(TEST_URLS['text_iso'])
        assert isinstance(content, str)
        assert content == "Hello, World!"

    def test_local_text_file_encodings(self, text_extractor):
        """Test handling of local text files with different encodings."""
        test_text = "Hello, World! 你好，世界！"

        # Test UTF-8
        with tempfile.NamedTemporaryFile(suffix='.txt', delete=False) as temp_txt:
            temp_txt.write(test_text.encode('utf-8'))
            temp_txt_path = temp_txt.name

        try:
            content = text_extractor.extract_content(temp_txt_path)
            assert content == test_text
        finally:
            os.unlink(temp_txt_path)

        # Test UTF-16
        with tempfile.NamedTemporaryFile(suffix='.txt', delete=False) as temp_txt:
            temp_txt.write(test_text.encode('utf-16'))
            temp_txt_path = temp_txt.name

        try:
            content = text_extractor.extract_content(temp_txt_path)
            assert content == test_text
        finally:
            os.unlink(temp_txt_path)

    def test_empty_text_file(self, text_extractor):
        """Test handling of empty text files."""
        with tempfile.NamedTemporaryFile(suffix='.txt', delete=False) as temp_txt:
            temp_txt_path = temp_txt.name

        try:
            content = text_extractor.extract_content(temp_txt_path)
            assert content == ""
        finally:
            os.unlink(temp_txt_path)

    @patch('requests.get')
    def test_invalid_url(self, mock_get, text_extractor):
        """Test handling of invalid URLs."""
        mock_get.return_value = create_mock_response(b'', status_code=404)
        with pytest.raises(requests.exceptions.HTTPError):
            text_extractor.extract_content("https://invalid.url/file.txt")

def test_encoding_handling(content_extractor):
    """Test handling of different text encodings."""
    # Test UTF-8 with BOM
    with tempfile.NamedTemporaryFile(suffix='.txt', delete=False) as temp_txt:
        temp_txt.write('Hello World'.encode('utf-8-sig'))
        temp_txt_path = temp_txt.name

    try:
        content = content_extractor.extract_content(temp_txt_path)
        assert content == 'Hello World'
    finally:
        os.unlink(temp_txt_path)

    # Test UTF-16
    with tempfile.NamedTemporaryFile(suffix='.txt', delete=False) as temp_txt:
        temp_txt.write('Hello World'.encode('utf-16'))
        temp_txt_path = temp_txt.name

    try:
        content = content_extractor.extract_content(temp_txt_path)
        assert 'Hello World' in content
    finally:
        os.unlink(temp_txt_path)

class TestPDFExtractor:
    @patch('requests.get')
    def test_web_pdf_extraction(self, mock_get, pdf_extractor):
        """Test PDF extraction from web URLs."""
        # Create a simple PDF content mock
        mock_pdf_content = b'%PDF-1.4\nsome content'
        mock_get.return_value = create_mock_response(mock_pdf_content)

        # Test successful extraction
        content = pdf_extractor.extract_content(TEST_URLS['pdf'])
        assert isinstance(content, str)
        mock_get.assert_called_once()

        # Test failed download
        mock_get.reset_mock()
        mock_get.return_value = create_mock_response(b'', status_code=404)
        with pytest.raises(requests.exceptions.HTTPError):
            pdf_extractor.extract_content(TEST_URLS['pdf'])

    def test_password_protected_pdf(self, pdf_extractor):
        """Test handling of password-protected PDFs."""
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
            # Create a password-protected PDF here
            temp_pdf.write(b'%PDF-1.4\n...encrypted...')
            temp_pdf_path = temp_pdf.name

        try:
            with pytest.raises(ValueError, match="PDF is password-protected"):
                pdf_extractor.extract_content(temp_pdf_path)
        finally:
            os.unlink(temp_pdf_path)

    def test_corrupted_pdf(self, pdf_extractor):
        """Test handling of corrupted PDFs."""
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
            temp_pdf.write(b'This is not a valid PDF file')
            temp_pdf_path = temp_pdf.name

        try:
            with pytest.raises(ValueError, match="PDF file is corrupted or invalid"):
                pdf_extractor.extract_content(temp_pdf_path)
        finally:
            os.unlink(temp_pdf_path)

    def test_empty_pdf(self, pdf_extractor):
        """Test handling of empty PDFs."""
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
            temp_pdf.write(b'%PDF-1.4\n%%EOF')  # Minimal valid PDF
            temp_pdf_path = temp_pdf.name

        try:
            content = pdf_extractor.extract_content(temp_pdf_path)
            assert content == ""
        finally:
            os.unlink(temp_pdf_path)

    def test_pdf_with_images(self, pdf_extractor):
        """Test handling PDFs with images (should only extract text)."""
        # Create a PDF with both text and images
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
            # A minimal PDF with text and an image placeholder
            pdf_content = b'''%PDF-1.4
1 0 obj
<</Type/Catalog/Pages 2 0 R>>
endobj
2 0 obj
<</Type/Pages/Kids[3 0 R]/Count 1>>
endobj
3 0 obj
<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]/Contents 4 0 R/Resources<</Font<</F1<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>>>>>
endobj
4 0 obj
<</Length 68>>
stream
BT
/F1 12 Tf
72 720 Td
(This is text that should be extracted) Tj
ET
endstream
endobj
xref
0 5
0000000000 65535 f
0000000009 00000 n
0000000056 00000 n
0000000107 00000 n
0000000247 00000 n
trailer
<</Size 5/Root 1 0 R>>
startxref
364
%%EOF'''
            temp_pdf.write(pdf_content)
            temp_pdf_path = temp_pdf.name

        try:
            content = pdf_extractor.extract_content(temp_pdf_path)
            assert "This is text that should be extracted" in content
        finally:
            os.unlink(temp_pdf_path)

    def test_pdf_with_special_characters(self, pdf_extractor):
        """Test handling PDFs with special characters and different encodings."""
        special_chars_text = "Hello, 世界! ¡Hola! Café"
        
        # Create a PDF with special characters
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
            # A minimal PDF with UTF-8 text
            pdf_content = f'''%PDF-1.4
1 0 obj
<</Type/Catalog/Pages 2 0 R>>
endobj
2 0 obj
<</Type/Pages/Kids[3 0 R]/Count 1>>
endobj
3 0 obj
<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]/Contents 4 0 R/Resources<</Font<</F1<</Type/Font/Subtype/Type1/BaseFont/Helvetica-UTF8>>>>>>
endobj
4 0 obj
<</Length {len(special_chars_text)}>>
stream
BT
/F1 12 Tf
72 720 Td
({special_chars_text}) Tj
ET
endstream
endobj
xref
0 5
0000000000 65535 f
0000000009 00000 n
0000000056 00000 n
0000000107 00000 n
0000000247 00000 n
trailer
<</Size 5/Root 1 0 R>>
startxref
364
%%EOF'''.encode('utf-8')
            temp_pdf.write(pdf_content)
            temp_pdf_path = temp_pdf.name

        try:
            content = pdf_extractor.extract_content(temp_pdf_path)
            # The extractor should handle the special characters appropriately
            assert isinstance(content, str)
            assert "Hello" in content
            # Note: Some special characters might be normalized or replaced
            # depending on the PDF encoding and extraction settings
        finally:
            os.unlink(temp_pdf_path)

    @patch('requests.get')
    def test_large_pdf_handling(self, mock_get, pdf_extractor):
        """Test handling of large PDFs."""
        # Create a large PDF content
        large_text = "Test content\n" * 10000  # Create a large amount of text
        pdf_content = f'''%PDF-1.4
1 0 obj
<</Type/Catalog/Pages 2 0 R>>
endobj
2 0 obj
<</Type/Pages/Kids[3 0 R]/Count 1>>
endobj
3 0 obj
<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]/Contents 4 0 R/Resources<</Font<</F1<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>>>>>
endobj
4 0 obj
<</Length {len(large_text)}>>
stream
BT
/F1 12 Tf
72 720 Td
({large_text}) Tj
ET
endstream
endobj
xref
0 5
0000000000 65535 f
0000000009 00000 n
0000000056 00000 n
0000000107 00000 n
0000000247 00000 n
trailer
<</Size 5/Root 1 0 R>>
startxref
364
%%EOF'''.encode('utf-8')

        mock_get.return_value = create_mock_response(pdf_content)
        
        # Test extraction of large PDF
        content = pdf_extractor.extract_content(TEST_URLS['pdf'])
        assert isinstance(content, str)
        assert len(content) > 50000  # Verify we got a large amount of text
        assert "Test content" in content

if __name__ == "__main__":
    unittest.main()
